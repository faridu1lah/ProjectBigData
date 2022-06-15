import streamlit as st
import base64
from pptx import Presentation
from io import BytesIO
import streamlit.components.v1 as components


def load_view():

    st.title("Download center for the developers!")
    # with st.form(key="my_form"):
      # Opening file from file path
    with open('RHerishIbrahim.pdf', "rb") as f:
        base64_pdf = base64.b64encode(f.read()).decode('utf-8')

    # Embedding PDF in HTML
    pdf_display = F'<iframe src="data:application/pdf;base64,{base64_pdf}" title="hallo" width="900" height="900" type="application/pdf">'

    # Displaying File
    st.markdown(pdf_display, unsafe_allow_html=True)

    prs = Presentation()
    # title_slide_layout = prs.slide_layouts[0]
    # slide = prs.slides.add_slide(title_slide_layout)
    # title = slide.shapes.title
    # subtitle = slide.placeholders[1]

    # title.text = "Hello, World!"
    # subtitle.text = "python-pptx was here!"

# save the output into binary form
    binary_output = BytesIO()
    prs.save(binary_output) 

    # st.download_button(label='Download ppw',
    #                     data=binary_output.getvalue(),
    #                     file_name='ADS.pptx') 

    st.write("## Download links:")
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.header("SPRINT 1 ")
        st.download_button(label='PowerPoint Sprint 1',data=binary_output.getvalue(),file_name='ADS.pptx') 

    with col2:
        st.header("SPRINT 2")
        st.download_button(label='Sprint 2',data=binary_output.getvalue(),file_name='ADS.pptx') 

    with col3:
        st.header("SPRINT 3")
        st.download_button(label='Sprint3 ppw',data=binary_output.getvalue(),file_name='ADS.pptx')    
    
    with col4:
        st.header("Source Code Link")
        st.write("[master.zip](https://github.com/FaridullahKhorshid/ProjectBigData/archive/refs/heads/master.zip)")